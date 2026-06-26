from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("outputs")
OUT.mkdir(exist_ok=True)

PPTX_PATH = OUT / "copyright_case_study_bartz_anthropic.pptx"
NOTES_PATH = OUT / "copyright_case_study_speaker_notes.md"
SOURCES_PATH = OUT / "copyright_case_study_sources.md"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(16, 31, 45)
BLUE = RGBColor(31, 92, 130)
TEAL = RGBColor(33, 145, 140)
AMBER = RGBColor(232, 170, 66)
RED = RGBColor(196, 72, 72)
PAPER = RGBColor(246, 247, 242)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(88, 98, 110)
LIGHT = RGBColor(229, 234, 238)


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, color=PAPER):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.z_order = 0


def add_text(slide, text, x, y, w, h, size=24, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None, section="CASE STUDY"):
    add_text(slide, section, 0.65, 0.35, 3.0, 0.3, 11, TEAL, True)
    add_text(slide, title, 0.65, 0.75, 10.9, 0.8, 32, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.68, 1.48, 10.9, 0.45, 15, MUTED)


def footer(slide, n, source="Sources: U.S. District Court order; Authors Guild; AP News; Clark Hill"):
    add_text(slide, source, 0.65, 7.08, 9.6, 0.25, 8, MUTED)
    add_text(slide, str(n), 12.45, 7.05, 0.3, 0.25, 9, MUTED)


def card(slide, x, y, w, h, title, body, accent=TEAL):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.25, y + 0.18, w - 0.45, 0.35, 17, NAVY, True)
    add_text(slide, body, x + 0.25, y + 0.62, w - 0.45, h - 0.75, 12.5, MUTED)


def bullet_list(slide, items, x, y, w, h, size=18, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(8)
    return box


# Slide 1
s = blank_slide()
add_bg(s, NAVY)
add_text(s, "Research Methodology Seminar", 0.75, 0.55, 5.5, 0.4, 14, AMBER, True)
add_text(s, "Copyright in the Age of AI", 0.75, 1.25, 8.2, 0.85, 40, WHITE, True)
add_text(s, "Case Study: Bartz v. Anthropic PBC and the $1.5B AI books settlement", 0.78, 2.15, 8.6, 0.55, 20, LIGHT)
add_text(s, "Focus: how a case study method explains a real legal-technology conflict", 0.78, 5.8, 8.6, 0.35, 15, LIGHT)
for x, y, c in [(9.5, 1.0, TEAL), (10.7, 2.3, AMBER), (9.7, 3.7, RED)]:
    circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.15), Inches(1.15))
    circ.fill.solid()
    circ.fill.fore_color.rgb = c
    circ.line.fill.background()
footer(s, 1, "Topic: Copyright case study | Latest status checked: June 2026")

# Slide 2
s = blank_slide()
add_bg(s)
add_title(s, "Why this case is interesting", "A current copyright dispute where authors, AI companies, markets, and law all collide.")
card(s, 0.8, 2.1, 3.7, 2.0, "New problem", "Generative AI systems need massive training data, including books and other creative works.", TEAL)
card(s, 4.8, 2.1, 3.7, 2.0, "Old legal question", "Copyright law asks whether copying is permitted, licensed, infringing, or protected by fair use.", AMBER)
card(s, 8.8, 2.1, 3.7, 2.0, "Research value", "The case gives clear documents, stakeholders, dates, evidence, and competing interpretations.", RED)
bullet_list(s, ["Suitable for a case study because it is bounded, recent, evidence-rich, and socially relevant."], 1.0, 5.1, 11.4, 0.7, 20)
footer(s, 2)

# Slide 3
s = blank_slide()
add_bg(s)
add_title(s, "Case background", "Bartz v. Anthropic PBC, U.S. District Court, Northern District of California.")
bullet_list(
    s,
    [
        "Filed by authors including Andrea Bartz, Charles Graeber, and Kirk Wallace Johnson.",
        "Claim: Anthropic copied books without permission while developing large language models.",
        "Key distinction: legally acquired books for AI training versus pirated books from shadow libraries.",
        "The case became a benchmark for AI training, licensing, and copyright risk.",
    ],
    0.9,
    2.0,
    7.5,
    3.1,
    18,
)
card(s, 9.0, 1.9, 3.3, 2.6, "Core tension", "Can an AI company learn from copyrighted books without permission, and does the source of those books change the answer?", BLUE)
footer(s, 3)

# Slide 4
s = blank_slide()
add_bg(s)
add_title(s, "Timeline of the case", "The case is useful because events can be arranged chronologically.")
events = [
    ("Aug 2024", "Authors sue Anthropic"),
    ("Jun 2025", "Court says AI training can be fair use when books are lawfully acquired"),
    ("Jul 2025", "Class certified for piracy-related claims"),
    ("Sep 2025", "$1.5B settlement receives preliminary approval"),
    ("Jun 2026", "Final approval still pending"),
]
x = 0.75
for i, (date, text) in enumerate(events):
    cx = x + i * 2.45
    line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(cx), Inches(3.1), Inches(2.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL if i < 4 else AMBER
    line.line.fill.background()
    circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx + 0.75), Inches(2.75), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL if i < 4 else AMBER
    circ.line.fill.background()
    add_text(s, date, cx, 3.45, 2.1, 0.35, 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(s, text, cx - 0.1, 3.85, 2.25, 0.9, 11.5, MUTED, False, PP_ALIGN.CENTER)
footer(s, 4, "Sources: Authors Guild update; AP News; Clark Hill June 2026 update")

# Slide 5
s = blank_slide()
add_bg(s)
add_title(s, "Research methodology design", "Single-case qualitative case study with document analysis.")
card(s, 0.8, 1.95, 3.7, 1.8, "Case type", "Instrumental case study: the case is used to understand a wider issue, AI and copyright.", TEAL)
card(s, 4.8, 1.95, 3.7, 1.8, "Data sources", "Court order, settlement notice, author advocacy updates, legal commentary, and news reports.", BLUE)
card(s, 8.8, 1.95, 3.7, 1.8, "Analysis method", "Thematic coding around fair use, piracy, market harm, licensing, and stakeholder impact.", AMBER)
card(s, 2.5, 4.55, 3.7, 1.45, "Unit of analysis", "The Bartz v. Anthropic dispute and its settlement pathway.", RED)
card(s, 6.9, 4.55, 3.7, 1.45, "Boundary", "U.S. copyright law context, 2024-2026 public record.", TEAL)
footer(s, 5)

# Slide 6
s = blank_slide()
add_bg(s)
add_title(s, "Research questions", "Questions guide the case study instead of simply narrating events.")
bullet_list(
    s,
    [
        "RQ1: How did the court distinguish AI training from unauthorized acquisition of copyrighted books?",
        "RQ2: What evidence shaped the fair-use and piracy analysis?",
        "RQ3: How did the settlement change incentives for authors, publishers, and AI companies?",
        "RQ4: What does the case suggest about future copyright licensing models for AI training data?",
    ],
    1.0,
    1.9,
    11.1,
    3.8,
    20,
)
footer(s, 6)

# Slide 7
s = blank_slide()
add_bg(s)
add_title(s, "Stakeholder map", "A case study identifies each actor and their interests.")
stakeholders = [
    ("Authors", "control, credit, compensation", TEAL, 1.0, 2.0),
    ("Publishers", "rights management, licensing markets", BLUE, 5.0, 1.4),
    ("Anthropic", "training data, AI product development", AMBER, 8.9, 2.0),
    ("Court", "fair use, class action fairness", RED, 5.0, 4.45),
]
for name, text, col, x, y in stakeholders:
    shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.0), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = col
    shape.line.width = Pt(2)
    add_text(s, name, x + 0.2, y + 0.18, 2.6, 0.25, 17, col, True, PP_ALIGN.CENTER)
    add_text(s, text, x + 0.2, y + 0.52, 2.6, 0.35, 10.5, MUTED, False, PP_ALIGN.CENTER)
add_text(s, "Copyrighted books + AI training data", 4.25, 3.05, 4.5, 0.45, 19, NAVY, True, PP_ALIGN.CENTER)
footer(s, 7)

# Slide 8
s = blank_slide()
add_bg(s)
add_title(s, "Key findings from the case", "The case does not give a simple yes/no answer.")
card(s, 0.8, 1.75, 3.7, 2.0, "Finding 1", "Training may be treated as transformative fair use when the material is lawfully acquired and not reproduced in outputs.", TEAL)
card(s, 4.8, 1.75, 3.7, 2.0, "Finding 2", "Pirated acquisition and retention created separate legal risk that fair use did not erase.", RED)
card(s, 8.8, 1.75, 3.7, 2.0, "Finding 3", "The settlement pushed the debate toward licensed datasets and cleaner provenance.", AMBER)
add_text(s, "Research interpretation", 0.9, 4.55, 3.0, 0.35, 19, NAVY, True)
bullet_list(s, ["The most important methodological lesson is to separate the legal issue into sub-questions: purpose of use, source of data, market effect, and remedy."], 0.95, 5.0, 11.0, 0.8, 18)
footer(s, 8)

# Slide 9
s = blank_slide()
add_bg(s)
add_title(s, "Quantitative snapshot", "Numbers help communicate the scale of the dispute.")
chart_data = CategoryChartData()
chart_data.categories = ["Eligible books", "Claimed works"]
chart_data.add_series("Works", (482460, 440490))
chart = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.9),
    Inches(1.9),
    Inches(6.5),
    Inches(3.7),
    chart_data,
).chart
chart.has_legend = False
chart.value_axis.has_major_gridlines = True
chart.category_axis.tick_labels.font.size = Pt(12)
chart.value_axis.tick_labels.font.size = Pt(11)
card(s, 8.0, 2.0, 4.2, 1.15, "$1.5B", "Proposed settlement amount plus interest", RED)
card(s, 8.0, 3.45, 4.2, 1.15, "about $3,000", "Estimated payment per qualifying work before deductions/splits", AMBER)
card(s, 8.0, 4.9, 4.2, 1.15, "91.3%", "Claim rate reported before the May 2026 fairness hearing", TEAL)
footer(s, 9, "Sources: Authors Guild; Society of Authors; AP News")

# Slide 10
s = blank_slide()
add_bg(s)
add_title(s, "Ethical issues in the study", "Research methodology also evaluates ethics, not just facts.")
bullet_list(
    s,
    [
        "Respect for creators: creative labor has economic and moral value.",
        "Data provenance: AI datasets should record where training material came from.",
        "Transparency: class members need clear notice, claim access, and fair payment procedures.",
        "Balance: innovation benefits society, but shortcuts in data collection can create harm.",
    ],
    1.0,
    1.85,
    10.8,
    3.5,
    20,
)
footer(s, 10)

# Slide 11
s = blank_slide()
add_bg(s)
add_title(s, "Conclusion", "Bartz v. Anthropic is a strong case study for copyright research methodology.")
card(s, 0.9, 1.85, 3.5, 2.2, "Method lesson", "A case study works best when the researcher defines clear boundaries, data sources, and research questions.", TEAL)
card(s, 4.9, 1.85, 3.5, 2.2, "Copyright lesson", "AI training, fair use, piracy, licensing, and market harm must be analyzed separately.", AMBER)
card(s, 8.9, 1.85, 3.5, 2.2, "Practical lesson", "Future AI companies may need stronger licensing, dataset audits, and creator compensation systems.", RED)
add_text(s, "Seminar closing line", 0.9, 5.05, 2.8, 0.35, 18, NAVY, True)
add_text(s, "This case shows that research methodology is not only about collecting facts; it is about organizing evidence so a complex real-world problem becomes understandable.", 0.9, 5.45, 11.2, 0.7, 18, MUTED)
footer(s, 11)

# Slide 12
s = blank_slide()
add_bg(s)
add_title(s, "References", "Use these sources in your final slide or seminar handout.")
refs = [
    "U.S. District Court order in Bartz v. Anthropic PBC, June 23, 2025.",
    "Authors Guild, 'Bartz v. Anthropic Settlement: What Authors Need to Know', updated Apr. 8, 2026.",
    "AP News, preliminary $1.5B settlement approval report, Sept. 2025.",
    "Clark Hill, 'Right To Know - June 2026, Vol. 42', June 15, 2026.",
    "Society of Authors, Anthropic settlement update, 2026.",
]
bullet_list(s, refs, 0.95, 1.95, 11.4, 4.0, 16)
footer(s, 12, "References are summarized for classroom use; verify legal advice with primary materials.")

prs.save(PPTX_PATH)

NOTES_PATH.write_text(
    """# Speaker Notes: Copyright Case Study Seminar

## Topic
Copyright in the Age of AI: Case Study of Bartz v. Anthropic PBC

## Opening
Good morning/afternoon. My seminar uses a case study method to examine a recent copyright dispute involving AI training data. The case is Bartz v. Anthropic PBC, involving authors, copyrighted books, AI model training, and a proposed $1.5 billion settlement.

## Why this case
This case is interesting because it is not only about law. It involves technology, creative labor, market value, research ethics, data collection, and corporate responsibility.

## Methodology explanation
I selected a single-case qualitative case study. The unit of analysis is the Bartz v. Anthropic dispute. The method is document analysis using court orders, settlement updates, author-organization materials, and news reports.

## Main points to say
- AI companies need large datasets, but copyright owners argue that their works should not be copied without permission.
- The court separated legally acquired books used for AI training from pirated copies obtained from shadow libraries.
- The case suggests that the source of data matters, not just the final use.
- The settlement is important because it may influence future licensing models for AI training data.

## Latest status
As of a June 15, 2026 legal update, final approval of the $1.5 billion settlement remained pending after the May 2026 fairness hearing.

## Closing
This case shows how research methodology helps organize a complex real-world issue. By defining the case boundary, research questions, data sources, and analysis themes, we can understand the legal and ethical debate more clearly.
""",
    encoding="utf-8",
)

SOURCES_PATH.write_text(
    """# Sources

1. U.S. District Court, Northern District of California, Bartz v. Anthropic PBC, Order on Fair Use, June 23, 2025.
   https://copyrightalliance.org/wp-content/uploads/2025/06/Bartz-v.-Anthropic-Order.pdf

2. Authors Guild, "Bartz v. Anthropic Settlement: What Authors Need to Know", updated April 8, 2026.
   https://authorsguild.org/advocacy/artificial-intelligence/what-authors-need-to-know-about-the-anthropic-settlement/

3. Associated Press, "Judge approves $1.5 billion copyright settlement between AI company Anthropic and authors", September 2025.
   https://apnews.com/article/anthropic-authors-copyright-judge-artificial-intelligence-9643064e847a5e88ef6ee8b620b3a44c

4. Clark Hill, "Right To Know - June 2026, Vol. 42", June 15, 2026.
   https://www.clarkhill.com/news-events/news/right-to-know-june-2026-vol-42/

5. Society of Authors, "The Anthropic settlement", January 20, 2026.
   https://societyofauthors.org/2026/01/20/anthropic-list-of-stolen-works-published/
""",
    encoding="utf-8",
)

print(PPTX_PATH)
print(NOTES_PATH)
print(SOURCES_PATH)
